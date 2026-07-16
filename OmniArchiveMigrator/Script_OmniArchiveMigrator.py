# ==========================================================
# MODULE:       Script_OmniArchiveMigrator
# PURPOSE:      自動掃描指定資料夾，判定檔案命名品質，提取內容特徵智慧重新命名，並安全複製至專屬目錄，同時生成溯源對照表
# EXPORTS:      UltimateArchaeologyGod
# IMPORTS:      os, shutil, re, csv, hashlib, pathlib, datetime, docx, openpyxl, pypdf, pptx
# FORBIDDEN:    禁止使用 open('w') 直接覆寫正式報表；禁止使用未經驗證的直接移動（shutil.move）
# DEPENDENCIES: python-docx, openpyxl, pypdf, python-pptx
# VERSION:      1.0.0 [Stability: Experimental]
# NOTES:        檔案特徵提取功能若遇首頁空白、全圖片或密碼保護之檔案，會預期性回傳「內容辨識失敗」，相關原因將記錄於報表備註欄。
# ==========================================================

import os, shutil, re, csv, hashlib
from pathlib import Path
from datetime import datetime

# ==========================================
# 考古庫加載 (Office/PDF/PPT)
# ==========================================
try:
    from docx import Document
    from openpyxl import load_workbook
    from pypdf import PdfReader
    from pptx import Presentation  # 需要執行: pip install python-pptx
    HAS_LIBS = True
except ImportError:
    HAS_LIBS = False

CONFIG = {
    "WS_PREFIX": "GOD_ULTIMATE_VAULT",
    "UGLY_DIR": "流浪漢隔離區",
    "BEAUTY_DIR": "美名保存區",
    "UGLY_PATTERNS": [r'新增', r'document', r'doc\d+', r'新建', r'text\d+', r'temp', r'无标题', r'copy of', r'presentation'],
    "REPORT": "00_全能歸檔診斷報告.csv"
}

def is_ugly(stem):
    """【改名判定】"""
    if stem.isdigit() or len(stem) < 3: return True
    if any(re.search(p, stem.lower()) for p in CONFIG["UGLY_PATTERNS"]): return True
    return False

def calculate_md5(file_path):
    """【特徵計算】計算檔案 MD5 用於去重與防碰撞"""
    hasher = hashlib.md5()
    try:
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hasher.update(chunk)
        return hasher.hexdigest()
    except Exception:
        return "ERROR_HASH"

class UltimateArchaeologyGod:
    def __init__(self, src_path):
        self.src = Path(src_path).resolve()
        ts = datetime.now().strftime('%m%d_%H%M')
        self.ws = self.src.parent / f"{CONFIG['WS_PREFIX']}_{ts}"
        self.beauty = self.ws / CONFIG["BEAUTY_DIR"]
        self.ugly = self.ws / CONFIG["UGLY_DIR"]
        for p in [self.beauty, self.ugly]: 
            p.mkdir(parents=True, exist_ok=True)
        self.log = []
        self.seen_hashes = {}  # 用於 ACDS 標準去重: {md5: 原始路徑}
        self.error_counter = 0 # 用於 ERROR_HASH 雙重失敗時的唯一性保底

    def _rescue_txt(self, path):
        """【萬用轉檔】僅用於提取標籤，不變更原始檔案內容"""
        err_log = []
        for enc in ['utf-8-sig', 'utf-8', 'gb18030', 'cp950', 'big5']:
            try:
                content = path.read_text(encoding=enc, errors='strict')
                if content.strip(): 
                    return content, enc, ""
            except Exception as e:
                err_log.append(f"{enc}:{str(e)}")
                continue
        
        err_msg = " | ".join(err_log)
        try:
            return path.read_text(encoding='utf-8', errors='ignore'), "FORCE_UTF8", f"嚴格解碼全敗({err_msg})"
        except Exception as e:
            return "", "DECODE_FAIL", f"終極讀取失敗: {str(e)} | 先前錯誤: {err_msg}"

    def _get_office_info(self, path):
        """【安全探針】提取關鍵字，並捕捉例外作為除錯線索"""
        ext = path.suffix.lower()
        if not HAS_LIBS: return "", "未安裝Office解析套件"
        
        try:
            if ext == ".pptx":
                prs = Presentation(path)
                text_runs = [shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text")]
                return " ".join(text_runs), ""
            elif ext == ".docx": 
                return " ".join([p.text for p in Document(path).paragraphs[:2]]), ""
            elif ext == ".xlsx": 
                wb = load_workbook(path, read_only=True, data_only=True)
                val = next(wb.active.values, "")
                wb.close()
                return str(val), ""
            elif ext == ".pdf": 
                return PdfReader(path).pages[0].extract_text()[:100], ""
        except Exception as e:
            return "", f"探針解析異常: {str(e)}"
        return "", ""

    def execute(self):
        exts = {'.txt', '.docx', '.xlsx', '.pdf', '.pptx'}
        files = [f for f in self.src.rglob("*") if f.suffix.lower() in exts]
        print(f"開始掃描與全方位歸檔中...")

        try:
            for i, src in enumerate(files, 1):
                if CONFIG["WS_PREFIX"] in str(src): continue
                
                # 1. 計算 MD5 與去重驗證
                file_md5 = calculate_md5(src)
                if file_md5 != "ERROR_HASH" and file_md5 in self.seen_hashes:
                    self.log.append([src.name, "N/A", file_md5, "N/A", "[重複跳過]", "成功", f"與 {self.seen_hashes[file_md5]} 內容相同"])
                    print(f"[{i:03d}] [重複跳過] -> {src.name[:25]}...")
                    continue
                
                # 建立安全檔案特徵碼 UID (保底處理雙重讀取錯誤)
                if file_md5 != "ERROR_HASH":
                    self.seen_hashes[file_md5] = src.name
                    file_uid = file_md5
                else:
                    self.error_counter += 1
                    file_uid = f"ERR{self.error_counter:04d}"

                # 2. 內容辨識 (僅作為命名依據)
                ext = src.suffix.lower()
                if ext == ".txt":
                    preview, enc, debug_msg = self._rescue_txt(src)
                else:
                    preview, debug_msg = self._get_office_info(src)
                    enc = f"{ext.upper()}_SAFE"

                # 3. 分流判定與智慧命名
                target = self.ugly if is_ugly(src.stem) else self.beauty
                tags = re.findall(r'[\u4e00-\u9fff]+', preview)
                tag_str = "".join(tags)[:8] if tags else "內容辨識失敗"
                date_str = datetime.fromtimestamp(src.stat().st_mtime).strftime('%Y%m%d')

                if target == self.ugly:
                    base_new_name = f"{date_str}_{tag_str}_{file_uid[:4]}"
                    action = "[隔離改名]"
                else:
                    base_new_name = src.stem
                    action = "[保存美名]"

                # 4. 終極防碰撞機制 (解決同名覆寫與雙重 ERROR_HASH 情境)
                final_new_name = f"{base_new_name}{ext}"
                dest = target / final_new_name
                bump = 0
                original_debug = debug_msg
                
                while dest.exists():
                    bump += 1
                    bump_suffix = f"_{file_uid[:6]}" + (f"_{bump}" if bump > 1 else "")
                    final_new_name = f"{base_new_name}{bump_suffix}{ext}"
                    dest = target / final_new_name
                    debug_msg = f"檔名碰撞已處理. {original_debug}".strip()

                # 5. 物理執行 (一律進行位元組級安全搬運)
                try:
                    shutil.copy2(src, dest)
                    print(f"[{i:03d}] {action} ({enc}) -> {final_new_name[:25]}...")
                    self.log.append([src.name, final_new_name, file_md5, enc, action, "成功", debug_msg])
                except Exception as e:
                    self.log.append([src.name, final_new_name, file_md5, enc, action, "失敗", str(e)])

        except KeyboardInterrupt:
            print("\n[中斷] 偵測到使用者取消 (Ctrl+C)，正在保留已處理的紀錄...")
            self.log.append(["系統中斷", "N/A", "N/A", "N/A", "Fail-Fast", "中斷", "使用者強制結束執行"])
        except Exception as e:
            print(f"\n[錯誤] 發生未預期錯誤，正在保留已處理的紀錄: {str(e)}")
            self.log.append(["系統錯誤", "N/A", "N/A", "N/A", "Exception", "失敗", str(e)])
        finally:
            self._save_report()
            self._finish_ui()

    def _save_report(self):
        """【報表保留】"""
        with open(self.ws / CONFIG["REPORT"], "w", newline="", encoding="utf-8-sig") as f:
            csv.writer(f).writerows([["原名", "新名", "MD5", "辨識編碼", "動作", "狀態", "備註"]] + self.log)

    def _finish_ui(self):
        """【停止保留】"""
        print("\n" + "="*60)
        print(f"歸檔完成！產出路徑: {self.ws.name}")
        print(f"獨立檔案數: {len(self.seen_hashes)}")
        print("*"*30)
        input("請在檢查報表後，按 Enter 鍵結束關閉...")

if __name__ == "__main__":
    p = input("請輸入要全方位歸檔的路徑: ").strip().strip('"')
    if os.path.isdir(p): 
        UltimateArchaeologyGod(p).execute()
    else:
        print("找不到該路徑，請確認後再試。")
