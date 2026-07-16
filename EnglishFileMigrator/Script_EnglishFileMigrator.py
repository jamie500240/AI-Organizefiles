# ==========================================================
# MODULE:       Script_EnglishFileMigrator
# PURPOSE:      自動掃描指定資料夾，判定語系並安全移轉英文檔案至專屬目錄，同時生成溯源對照表
# EXPORTS:      EnglishFileMigrator
# IMPORTS:      os, shutil, csv, hashlib, logging, datetime, pathlib, docx, pdfminer, pptx, openpyxl, pytesseract, PIL
# FORBIDDEN:    禁止使用 open('w') 直接覆寫正式報表；禁止使用未經驗證的直接移動（shutil.move）
# DEPENDENCIES: ACDS_ContractRegistry (CONFIG 變數結構), ACDS_ADR (ADR-001, ADR-002, ADR-003, ADR-004)
# VERSION:      2.0.0 [Stability: Experimental]
# ==========================================================

import sys
import os
import shutil
import csv
import hashlib
import logging
from datetime import datetime
from pathlib import Path

# 解析套件
from docx import Document
from pdfminer.high_level import extract_text
from pptx import Presentation
import openpyxl

# 嘗試載入 OCR 相關套件
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# ==========================================
# SSOT: 語系自動判定配置 (唯一入口)
# ==========================================
CONFIG = {
    "CHINESE_THRESHOLD": 10,           # 絕對閥值：適用於極短字串（如僅檔名）
    "CHINESE_RATIO_THRESHOLD": 0.05,   # 比例閥值：長篇文本中文字元占比需小於 5%
    "PDF_PAGE_LIMIT": 30,
    "TARGET_EXTS": ['.pdf', '.docx', '.txt', '.xlsx', '.xls', '.csv', '.pptx',
                    '.jpg', '.jpeg', '.png', '.gif'],
    "REPORT_NAME": "英文檔案移轉對照表.csv",
    "ENGLISH_FOLDER": "英文區"
}

class ContentExtractionError(Exception):
    """自訂例外：內容解析失敗，避免靜默回傳 False"""
    pass

class EnglishFileMigrator:
    def __init__(self, root_dir, enable_ocr=False, include_images_in_migration=False):
        self.root = Path(root_dir).resolve()
        self.eng_folder = self.root / CONFIG["ENGLISH_FOLDER"]
        self.report_path = self.root / CONFIG["REPORT_NAME"]
        self.temp_report_path = self.root / f"{CONFIG['REPORT_NAME']}.tmp"
        self.stats = {"moved": 0, "stayed": 0, "errors": 0}
        
        # 參數設定
        self.enable_ocr = enable_ocr and OCR_AVAILABLE
        self.include_images = include_images_in_migration
        
        # 初始化 Logging
        self.logger = logging.getLogger("EnglishMigrator")
        self.logger.setLevel(logging.INFO)
        if not self.logger.handlers:
            ch = logging.StreamHandler()
            formatter = logging.Formatter('[%(levelname)s] %(asctime)s 英文移轉模組: %(message)s', datefmt='%Y-%m-%d %H:%M:%S')
            ch.setFormatter(formatter)
            self.logger.addHandler(ch)

        if enable_ocr and not OCR_AVAILABLE:
            self.logger.warning("未偵測到 pytesseract 或 Pillow，OCR 功能已被強制關閉。請確認是否已安裝套件與 Tesseract-OCR。")

    def _get_file_hash(self, path):
        """計算檔案 SHA-256 雜湊值"""
        hasher = hashlib.sha256()
        try:
            with open(path, 'rb') as f:
                for chunk in iter(lambda: f.read(65536), b''):
                    hasher.update(chunk)
            return hasher.hexdigest()
        except PermissionError:
            raise
        except Exception as e:
            self.logger.error(f"計算雜湊失敗 {path.name}: {e}")
            return None

    # ---------------------------------------------------------
    # 格式專屬解析器 (Extractors)
    # ---------------------------------------------------------
    def _extract_pdf(self, path):
        return extract_text(path, page_numbers=list(range(CONFIG['PDF_PAGE_LIMIT'])))

    def _extract_docx(self, path):
        doc = Document(path)
        return " ".join([p.text for i, p in enumerate(doc.paragraphs) if i < CONFIG['PDF_PAGE_LIMIT']])

    def _extract_txt(self, path):
        try:
            return path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            return path.read_text(encoding='big5', errors='ignore')

    def _extract_csv(self, path):
        lines = []
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for _, row in zip(range(10), reader):
                lines.append(" ".join(row))
        return " ".join(lines)

    def _extract_xlsx(self, path):
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheet = wb.active
        lines = []
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            if i >= 10: break
            lines.append(" ".join([str(cell) for cell in row if cell is not None]))
        wb.close()
        return " ".join(lines)

    def _extract_pptx(self, path):
        prs = Presentation(path)
        text_parts = []
        for slide in prs.slides[:CONFIG['PDF_PAGE_LIMIT']]:
            text_parts.extend([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        return " ".join(text_parts)

    def _extract_image(self, path):
        try:
            img = Image.open(path)
            # lang='eng+chi_tra' 支援中英混合辨識
            return pytesseract.image_to_string(img, lang='eng+chi_tra') 
        except Exception as e:
            self.logger.warning(f"圖片 OCR 辨識失敗 {path.name}: {e}")
            return ""

    # ---------------------------------------------------------
    # 核心語系判定邏輯
    # ---------------------------------------------------------
    def is_english_content(self, path):
        text_parts = [path.stem]
        ext = path.suffix.lower()
        
        try:
            if ext == '.pdf': text_parts.append(self._extract_pdf(path))
            elif ext == '.docx': text_parts.append(self._extract_docx(path))
            elif ext == '.txt': text_parts.append(self._extract_txt(path))
            elif ext == '.csv': text_parts.append(self._extract_csv(path))
            elif ext in ['.xlsx', '.xls']: text_parts.append(self._extract_xlsx(path))
            elif ext == '.pptx': text_parts.append(self._extract_pptx(path))
            elif ext in ['.jpg', '.jpeg', '.png', '.gif']:
                # 拒絕單憑檔名判斷圖片。若未啟用 OCR，一律拋出特定異常阻擋移轉
                if not self.enable_ocr:
                    raise ContentExtractionError("未啟用 OCR，僅供保留原位")
                text_parts.append(self._extract_image(path))
        except PermissionError:
            raise
        except ContentExtractionError:
            raise
        except Exception as e:
            raise ContentExtractionError(str(e))

        # 組合所有提取的字串
        full_text = " ".join(text_parts)
        text_length = len(full_text.strip())
        
        if text_length == 0:
            return False

        chinese_count = sum(1 for c in full_text if '\u4e00' <= c <= '\u9fff')
        
        if text_length < 50:
            return chinese_count < CONFIG['CHINESE_THRESHOLD']
        else:
            return (chinese_count / text_length) < CONFIG['CHINESE_RATIO_THRESHOLD']

    # ---------------------------------------------------------
    # 主執行流程
    # ---------------------------------------------------------
    def execute(self):
        if not self.root.is_dir():
            self.logger.critical(f"目標路徑無效: {self.root}")
            return

        self.eng_folder.mkdir(exist_ok=True)
        
        try:
            with open(self.temp_report_path, 'w', newline='', encoding='utf-8-sig') as f:
                writer = csv.writer(f)
                writer.writerow(["檔案名稱", "判定結果", "動作狀態", "原始位置(溯源用)", "目前位置"])

                for path in self.root.rglob("*"):
                    if not path.is_file(): continue
                    if path.suffix.lower() not in CONFIG["TARGET_EXTS"]: continue
                    if CONFIG["ENGLISH_FOLDER"] in path.parts: continue
                    if path.name == CONFIG["REPORT_NAME"] or path.name == self.temp_report_path.name: continue

                    old_full_path = str(path)
                    
                    try:
                        is_eng = self.is_english_content(path)
                    except PermissionError:
                        self.logger.warning(f"檔案遭鎖定，放棄讀取判定: {path.name}")
                        writer.writerow([path.name, "未判定", "檔案鎖定，保留原位", old_full_path, old_full_path])
                        self.stats["stayed"] += 1
                        continue
                    except ContentExtractionError as e:
                        # 處理 OCR 未啟動的特定標記狀態
                        action_str = "內容解析失敗，保留原位"
                        if str(e) == "未啟用 OCR，僅供保留原位":
                            action_str = str(e)
                        else:
                            self.logger.warning(f"檔案內文提取失敗，放棄判定: {path.name} (原因: {e})")
                            
                        writer.writerow([path.name, "未判定", action_str, old_full_path, old_full_path])
                        self.stats["stayed"] += 1
                        self.stats["errors"] += 1
                        continue
                    
                    # 圖片排外邏輯：除非明確允許，否則拒絕移轉圖片
                    if is_eng and (self.include_images or path.suffix.lower() not in ['.jpg', '.jpeg', '.png', '.gif']):
                        dest = self.eng_folder / path.name
                        
                        # 處理重複檔案與移轉
                        if dest.exists():
                            try:
                                src_hash = self._get_file_hash(path)
                                dest_hash = self._get_file_hash(dest)
                            except PermissionError:
                                self.logger.warning(f"檔案遭鎖定，放棄雜湊比對: {path.name}")
                                writer.writerow([path.name, "英文", "檔案鎖定，保留原位", old_full_path, old_full_path])
                                self.stats["stayed"] += 1
                                continue
                            
                            if src_hash is not None and src_hash == dest_hash:
                                action = "重複檔案（雜湊相同），已安全清理"
                                current_pos = str(dest)
                                try:
                                    path.unlink()
                                    self.stats["stayed"] += 1
                                except PermissionError:
                                    action = "清理重複檔案失敗: 權限不足或檔案佔用"
                                    self.logger.warning(action)
                                    self.stats["errors"] += 1
                                writer.writerow([path.name, "英文", action, old_full_path, current_pos])
                                continue
                            else:
                                counter = 1
                                while dest.exists():
                                    dest = self.eng_folder / f"{path.stem}_{counter}{path.suffix}"
                                    counter += 1

                        # 核心移轉邏輯 (Copy + Verify + Delete)
                        try:
                            shutil.copy2(path, dest)
                            
                            h1 = self._get_file_hash(path)
                            h2 = self._get_file_hash(dest)
                            
                            if h1 is not None and h1 == h2:
                                path.unlink()
                                action = "已移轉至英文區"
                                current_pos = str(dest)
                                self.stats["moved"] += 1
                                self.logger.info(f"成功安全移轉檔案: {path.name}")
                            else:
                                raise IOError("複製後雜湊值不一致或無法驗證，已清除殘留副本")
                                
                        except PermissionError:
                            action = "移轉失敗: 檔案遭系統鎖定 (In-Use)"
                            current_pos = old_full_path
                            self.stats["errors"] += 1
                            self.logger.warning(f"檔案佔用防禦觸發: {path.name}")
                            if dest.exists() and str(dest) != str(path):
                                try: 
                                    dest.unlink(missing_ok=True)
                                except Exception as clean_err: 
                                    self.logger.warning(f"清理殘留檔失敗 {dest.name}: {clean_err}")
                                
                        except Exception as e:
                            action = f"安全移轉失敗: {e}"
                            current_pos = old_full_path
                            self.stats["errors"] += 1
                            self.logger.error(f"移轉異常 {path.name}: {e}")
                            if dest.exists() and str(dest) != str(path):
                                try: 
                                    dest.unlink(missing_ok=True)
                                except Exception as clean_err: 
                                    self.logger.warning(f"清理異常殘留檔失敗 {dest.name}: {clean_err}")
                    else:
                        action = "保留原位"
                        current_pos = old_full_path
                        self.stats["stayed"] += 1
                        
                        msg = f"檔案保留原位: {path.name}"
                        # 補充日誌說明以便除錯
                        if path.suffix.lower() in ['.jpg', '.jpeg', '.png', '.gif'] and not self.include_images and is_eng:
                            msg = f"檔案保留原位(圖片安全排除規則): {path.name}"
                        elif not is_eng:
                            msg = f"檔案保留原位(非英文或判定未達標): {path.name}"
                            
                        self.logger.info(msg)

                    writer.writerow([path.name, "英文" if is_eng else "非英文", action, old_full_path, current_pos])

            self.temp_report_path.replace(self.report_path)
            self._show_summary()

        except (Exception, KeyboardInterrupt) as e:
            err_msg = "使用者強制中斷" if isinstance(e, KeyboardInterrupt) else f"批次任務崩潰: {e}"
            self.logger.critical(err_msg)
            
            if self.temp_report_path.exists():
                crash_log_path = self.root / f"{CONFIG['REPORT_NAME']}.crashed"
                try:
                    self.temp_report_path.rename(crash_log_path)
                    self.logger.critical(f"為保留溯源軌跡，未完成之報表已另存為: {crash_log_path.name}")
                except Exception as rename_err:
                    self.logger.critical(f"無法保留崩潰報表: {rename_err}")
                    
            if isinstance(e, KeyboardInterrupt):
                raise  

    def _show_summary(self):
        self.logger.info("-" * 50)
        self.logger.info("移轉任務執行結束")
        self.logger.info(f"統計: 已移動 {self.stats['moved']} | 保留原位 {self.stats['stayed']} | 異常 {self.stats['errors']}")
        self.logger.info("（註：發生異常的檔案皆已安全保留原位，因此「保留原位」包含了「異常」的數量）")
        self.logger.info(f"溯源報表路徑: {self.report_path}")
        self.logger.info("-" * 50)


if __name__ == "__main__":
    os.system('cls' if os.name == 'nt' else 'clear')
    
    print("="*60)
    print(" 英文檔案自動化移轉系統啟動 v2.0.1 [Experimental]")
    print("="*60)
    
    try:
        path_input = input("請輸入目標資料夾路徑: ").strip().strip('"')
        if not path_input:
            sys.exit()
            
        ocr_input = input("是否啟用圖片內容 OCR 辨識？(需安裝 tesseract) [y/N]: ").strip().lower()
        use_ocr = ocr_input in ['y', 'yes']
        
        img_input = input("是否允許移轉圖片檔案？(預設不搬移，以保護專案相簿與資料夾結構) [y/N]: ").strip().lower()
        move_imgs = img_input in ['y', 'yes']
            
        print("\n")
        migrator = EnglishFileMigrator(
            path_input, 
            enable_ocr=use_ocr,
            include_images_in_migration=move_imgs
        )
        migrator.execute()
    except KeyboardInterrupt:
        print("\n[WARN] 使用者強制中斷主程式")
    
    print("\n")
    input("任務結束，按 Enter 鍵關閉視窗...")
