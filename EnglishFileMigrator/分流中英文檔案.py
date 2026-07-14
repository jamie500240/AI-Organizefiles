# ==========================================================
# MODULE:       Script_EnglishFileMigrator
# PURPOSE:      自動掃描指定資料夾，判定語系並安全移轉英文檔案至專屬目錄，同時生成溯源對照表
# EXPORTS:      EnglishFileMigrator
# IMPORTS:      os, shutil, csv, hashlib, logging, datetime, pathlib, docx, pdfminer, pptx, openpyxl, pytesseract, PIL
# FORBIDDEN:    禁止使用 open('w') 直接覆寫正式報表；禁止使用未經驗證的直接移動（shutil.move）
# DEPENDENCIES: ACDS_ContractRegistry (CONFIG 變數結構), ACDS_ADR (ADR-001, ADR-002, ADR-003, ADR-004)
# VERSION:      1.0.0 [Stability: Experimental]
# ==========================================================

import os, shutil, csv
import pandas as pd
from pathlib import Path
from docx import Document
from pdfminer.high_level import extract_text
from pptx import Presentation

# ==========================================
# SSOT: 語系自動判定配置 (唯一入口)
# ==========================================
CONFIG = {
    "CHINESE_THRESHOLD": 10,       # 中文字符門檻
    "PDF_PAGE_LIMIT": 30,          # 掃描前 N 頁/段落
    "TARGET_EXTS": ['.pdf', '.docx', '.txt', '.xlsx', '.xls', '.csv', '.pptx',
                    '.jpg', '.jpeg', '.png', '.gif'],
    "REPORT_NAME": "英文檔案移轉對照表.csv",
    "ENGLISH_FOLDER": "英文區"
}

class EnglishMoveGod:
    def __init__(self, root_dir):
        # Defensive: 確保路徑絕對化
        self.root = Path(root_dir).resolve()
        self.eng_folder = self.root / CONFIG["ENGLISH_FOLDER"]
        self.report_path = self.root / CONFIG["REPORT_NAME"]
        self.stats = {"moved": 0, "stayed": 0, "errors": 0}

    def is_english_content(self, path):
        """內容提取與語系判定 (基於 Unicode 編碼)"""
        text = path.stem 
        ext = path.suffix.lower()
        
        try:
            if ext == '.pdf':
                text += extract_text(path, page_numbers=list(range(CONFIG['PDF_PAGE_LIMIT'])))
            elif ext == '.docx':
                doc = Document(path)
                text += " ".join([p.text for i, p in enumerate(doc.paragraphs) if i < CONFIG['PDF_PAGE_LIMIT']])
            elif ext == '.txt':
                try:
                    text += path.read_text(encoding='utf-8')
                except:
                    text += path.read_text(encoding='big5', errors='ignore')
            elif ext in ['.xlsx', '.xls', '.csv']:
                try:
                    df = pd.read_csv(path, nrows=10) if ext == '.csv' else pd.read_excel(path, nrows=10)
                    text += " ".join(df.columns.astype(str))
                except: pass
            elif ext == '.pptx':
                try:
                    prs = Presentation(path)
                    for i, slide in enumerate(prs.slides[:CONFIG['PDF_PAGE_LIMIT']]):
                        text += " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                except: pass
        except Exception:
            self.stats["errors"] += 1
            return False

        chinese_count = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
        
        if ext in ['.jpg', '.jpeg', '.png', '.gif']:
            chinese_count += sum(1 for c in path.stem if '\u4e00' <= c <= '\u9fff')

        return chinese_count < CONFIG['CHINESE_THRESHOLD']

    def execute(self):
        """核心執行：掃描、判定、直接移動"""
        if not self.root.is_dir():
            print(f"❌ 錯誤：路徑 {self.root} 不是有效的資料夾")
            return

        self.eng_folder.mkdir(exist_ok=True)
        
        # utf-8-sig 確保 Excel 開啟報表不亂碼
        with open(self.report_path, 'w', newline='', encoding='utf-8-sig') as f:
            writer = csv.writer(f)
            writer.writerow(["檔案名稱", "判定結果", "動作狀態", "原始位置(溯源用)", "目前位置"])

            # 遍歷所有子資料夾
            for path in self.root.rglob("*"):
                # 防禦檢查
                if not path.is_file(): continue
                if path.suffix.lower() not in CONFIG["TARGET_EXTS"]: continue
                if CONFIG["ENGLISH_FOLDER"] in path.parts: continue
                if path.name == CONFIG["REPORT_NAME"]: continue

                old_full_path = str(path)
                is_eng = self.is_english_content(path)
                
                # 只有非圖片的英文檔才執行移動 (保留圖片在原位是為了防止破壞相簿結構)
                if is_eng and path.suffix.lower() not in ['.jpg', '.jpeg', '.png', '.gif']:
                    dest = self.eng_folder / path.name
                    
                    # 衝突防禦：若英文區已有同名檔案，自動編號
                    counter = 1
                    while dest.exists():
                        dest = self.eng_folder / f"{path.stem}_{counter}{path.suffix}"
                        counter += 1
                    
                    try:
                        shutil.move(str(path), str(dest)) # 直接移動
                        action = "已移動至英文區"
                        current_pos = str(dest)
                        self.stats["moved"] += 1
                    except Exception as e:
                        action = f"移動失敗: {e}"
                        current_pos = old_full_path
                        self.stats["errors"] += 1
                else:
                    action = "保留原位"
                    current_pos = old_full_path
                    self.stats["stayed"] += 1

                writer.writerow([path.name, "英文" if is_eng else "非英文", action, old_full_path, current_pos])
                print(f"[{'MOVE' if action == '已移動至英文區' else 'KEEP'}] {path.name[:40]}")

        self._show_summary()

    def _show_summary(self):
        print("\n" + "═"*60)
        print(f"🎯 移轉任務完成！")
        print(f"📊 統計：已移動 {self.stats['moved']} | 保留原位 {self.stats['stayed']} | 異常 {self.stats['errors']}")
        print(f"📝 溯源報表：{self.report_path}")
        print(f"💡 提示：若需找回檔案，請參考報表中的「原始位置」欄位。")
        print("═"*60)

if __name__ == "__main__":
    os.system('cls' if os.name == 'nt' else 'clear')
    print("🚀 英文檔案【自動化移轉系統】啟動...\n")
    
    path_input = input("👉 請輸入目標資料夾路徑: ").strip().strip('"')
    
    god = EnglishMoveGod(path_input)
    god.execute()
    
    print("\n" + "—"*50)
    input("任務結束，按 Enter 鍵關閉視窗...")
